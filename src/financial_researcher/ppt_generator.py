
"""
Module to convert markdown report to PowerPoint presentation with professional design
"""

import os
import re
import urllib.request
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class ReportToPPT:
    """Convert markdown report to PowerPoint presentation with professional design"""
    
    # Professional color scheme
    PRIMARY_COLOR = RGBColor(13, 71, 161)      # Deep blue
    SECONDARY_COLOR = RGBColor(25, 103, 210)   # Vibrant blue
    ACCENT_COLOR = RGBColor(255, 152, 0)       # Amber/Orange accent
    TEXT_DARK = RGBColor(33, 33, 33)           # Dark gray
    TEXT_LIGHT = RGBColor(255, 255, 255)       # White
    BACKGROUND_LIGHT = RGBColor(245, 245, 245) # Light gray background
    
    def __init__(self, markdown_file: str, output_file: str = None, company_name: str = None):
        """
        Initialize the converter
        
        Args:
            markdown_file: Path to the markdown report file
            output_file: Path where the PowerPoint file will be saved
            company_name: Name of the company for logo lookup
        """
        self.markdown_file = markdown_file
        self.output_file = output_file or 'output/report.pptx'
        self.company_name = company_name
        self.company_logo_url = None
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Get company logo URL if company name is provided
        if self.company_name:
            self.company_logo_url = self._get_company_logo_url(self.company_name)
    
    def _get_company_logo_url(self, company_name: str) -> str:
        """
        Get the logo URL for a company by searching for it.
        Uses clearbit API to find company logo.
        
        Args:
            company_name: Name of the company
            
        Returns:
            URL string of the company logo or None if not found
        """
        try:
            # Using Clearbit API for company logo lookup
            # Convert company name to domain-friendly format
            domain = company_name.lower().replace(' ', '')
            logo_url = f"https://logo.clearbit.com/{domain}.com"
            
            # Verify the URL exists by making a HEAD request
            try:
                req = urllib.request.Request(logo_url, method='HEAD')
                with urllib.request.urlopen(req, timeout=5) as response:
                    if response.status == 200:
                        print(f"✓ Found company logo: {logo_url}")
                        return logo_url
            except:
                pass
            
            # Fallback: try with different domain extensions
            for ext in ['io', 'co', 'org', 'net']:
                logo_url = f"https://logo.clearbit.com/{domain}.{ext}"
                try:
                    req = urllib.request.Request(logo_url, method='HEAD')
                    with urllib.request.urlopen(req, timeout=5) as response:
                        if response.status == 200:
                            print(f"✓ Found company logo: {logo_url}")
                            return logo_url
                except:
                    continue
            
            print(f"⚠ Logo not found for {company_name}")
            return None
        except Exception as e:
            print(f"Warning: Could not fetch logo URL for {company_name}: {e}")
            return None
    
    
    def read_markdown(self) -> str:
        """Read the markdown file"""
        with open(self.markdown_file, 'r') as f:
            return f.read()
    
    def parse_markdown(self, content: str) -> list:
        """
        Parse markdown content into sections
        
        Returns:
            List of sections with title and content
        """
        sections = []
        current_section = {'title': '', 'content': []}
        
        lines = content.split('\n')
        
        for line in lines:
            # Main title (H1 or bold text)
            if line.startswith('# ') or (line.startswith('**') and line.endswith('**')):
                if current_section['title']:
                    sections.append(current_section)
                current_section = {'title': line.lstrip('# ').strip('*'), 'content': []}
            # Subsection (H2, H3, H4)
            elif re.match(r'^#{2,4}\s', line):
                if current_section['content'] or current_section['title']:
                    if current_section['title']:
                        sections.append(current_section)
                    current_section = {'title': line.lstrip('#').strip(), 'content': []}
            # Content
            elif line.strip():
                current_section['content'].append(line)
        
        if current_section['title']:
            sections.append(current_section)
        
        return sections
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add a professional title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add gradient background using shapes
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.PRIMARY_COLOR
        
        # Add accent shape on the side
        accent_shape = slide.shapes.add_shape(1, Inches(8.5), Inches(0), Inches(1.5), Inches(7.5))
        accent_fill = accent_shape.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = self.SECONDARY_COLOR
        accent_shape.line.color.rgb = self.SECONDARY_COLOR
        
        # Add top accent bar
        top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.2))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.ACCENT_COLOR
        top_bar.line.color.rgb = self.ACCENT_COLOR
        
        # Add company logo if available
        if self.company_logo_url:
            try:
                # Download logo temporarily
                logo_path = 'temp_logo.png'
                urllib.request.urlretrieve(self.company_logo_url, logo_path)
                # Add logo to top right of title slide
                slide.shapes.add_picture(logo_path, Inches(7.5), Inches(0.5), height=Inches(1.2))
                # Clean up temp file
                if os.path.exists(logo_path):
                    os.remove(logo_path)
                print(f"✓ Added logo to presentation")
            except Exception as e:
                print(f"Warning: Could not add logo image to slide: {e}")
        
        # Add title
        left = Inches(0.5)
        top = Inches(2.2)
        width = Inches(8)
        height = Inches(2)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_LIGHT
        p.alignment = PP_ALIGN.LEFT
        
        # Add subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(left, Inches(4.5), width, Inches(2))
            sub_frame = sub_box.text_frame
            sub_frame.word_wrap = True
            p = sub_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self.ACCENT_COLOR
            p.alignment = PP_ALIGN.LEFT
            p.font.italic = True
            
            # Add logo URL as text if company name is provided
            if self.company_name and self.company_logo_url:
                logo_text_box = slide.shapes.add_textbox(left, Inches(6.2), Inches(7), Inches(1))
                logo_text_frame = logo_text_box.text_frame
                logo_text_frame.word_wrap = True
                p = logo_text_frame.paragraphs[0]
                p.text = f"Logo Source: {self.company_logo_url}"
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(150, 150, 150)
                p.font.italic = True
    
    def add_content_slide(self, title: str, content_lines: list):
        """Add a professional content slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add top colored bar
        top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.ACCENT_COLOR
        top_bar.line.color.rgb = self.ACCENT_COLOR
        
        # Add left accent bar
        left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0.12), Inches(0.08), Inches(7.38))
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = self.PRIMARY_COLOR
        left_bar.line.color.rgb = self.PRIMARY_COLOR
        
        # Add title with background
        title_bg = slide.shapes.add_shape(1, Inches(0.1), Inches(0.3), Inches(9.8), Inches(0.8))
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.PRIMARY_COLOR
        title_bg.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_LIGHT
        p.alignment = PP_ALIGN.LEFT
        
        # Add content with better formatting
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Process content
        processed_content = []
        for line in content_lines:
            clean = line.strip().lstrip('- ').lstrip('* ').strip()
            if clean:
                processed_content.append(clean)
        
        # Add paragraphs with proper formatting
        for i, line in enumerate(processed_content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line
            p.font.name = 'Calibri'
            p.font.size = Pt(20)
            p.font.color.rgb = self.TEXT_DARK
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.line_spacing = 1.2
            
            # Add bullet points
            p.font.bold = False
    
    def add_summary_slide(self, title: str, key_points: list):
        """Add a summary/key points slide with visual emphasis"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Top colored bar
        top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.ACCENT_COLOR
        top_bar.line.color.rgb = self.ACCENT_COLOR
        
        # Left accent bar
        left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0.12), Inches(0.08), Inches(7.38))
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = self.PRIMARY_COLOR
        left_bar.line.color.rgb = self.PRIMARY_COLOR
        
        # Title
        title_bg = slide.shapes.add_shape(1, Inches(0.1), Inches(0.3), Inches(9.8), Inches(0.8))
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.PRIMARY_COLOR
        title_bg.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_LIGHT
        
        # Add key points
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(key_points):
            if not point.strip():
                continue
            
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            clean_point = point.strip().lstrip('- ').lstrip('* ').strip()
            p.text = clean_point
            p.font.name = 'Calibri'
            p.font.size = Pt(19)
            p.font.color.rgb = self.TEXT_DARK
            p.space_before = Pt(10)
            p.space_after = Pt(10)
            p.line_spacing = 1.25
            
            # Bullet point styling
            p.level = 0
    
    def generate(self) -> str:
        """
        Generate the PowerPoint presentation
        
        Returns:
            Path to the generated PowerPoint file
        """
        # Read and parse markdown
        content = self.read_markdown()
        sections = self.parse_markdown(content)
        
        # Extract title for first slide
        title = "Financial Research Report"
        subtitle = "Comprehensive Company Analysis & Market Outlook"
        
        # Check if first section looks like a title
        if sections and sections[0]['title']:
            title = sections[0]['title']
            sections = sections[1:]  # Remove title section
        
        # Add title slide
        self.add_title_slide(title, subtitle)
        
        # Add content slides
        for i, section in enumerate(sections):
            if section['title'].strip():
                # Alternate slide styles for visual variety
                if i % 3 == 2:  # Every 3rd slide can be a summary style
                    self.add_summary_slide(section['title'].strip(), section['content'])
                else:
                    self.add_content_slide(section['title'].strip(), section['content'])
        
        # Ensure output directory exists
        os.makedirs(os.path.dirname(self.output_file) or '.', exist_ok=True)
        
        # Save presentation
        self.prs.save(self.output_file)
        
        return self.output_file


def convert_report_to_ppt(markdown_path: str, ppt_path: str = None, company_name: str = None) -> str:
    """
    Convenience function to convert markdown report to PowerPoint
    
    Args:
        markdown_path: Path to the markdown report
        ppt_path: Path where the PowerPoint should be saved
        company_name: Name of the company for logo lookup
    
    Returns:
        Path to the generated PowerPoint file
    """
    converter = ReportToPPT(markdown_path, ppt_path, company_name=company_name)
    return converter.generate()
